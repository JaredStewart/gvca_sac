"""Auto-generate PPTX presentations from survey data.

Reads a template PPTX, replaces chart images and data-driven text on
data-driven slides, stamps narrative slides with a "NEEDS UPDATING" banner,
and writes the result to artifacts/{year}/presentation_{year}.pptx.
"""

import logging
import re
import tempfile
from pathlib import Path
from typing import Any

import polars as pl
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

from app.config import get_settings
from app.core.charts import (
    compute_level_average,
    compute_overall_average,
    compute_satisfaction_percentages,
    generate_category_yoy_chart,
    generate_composite_satisfaction_chart,
    generate_good_choice_better_serve_chart,
    generate_question_stacked_bar,
    generate_question_yoy_chart,
    generate_tag_frequency_chart,
)
from app.core.survey_config import (
    FREE_RESPONSE_QUESTIONS,
    LEVELS,
    QUESTION_SCALES,
    QUESTIONS,
)

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Slide classification
# ---------------------------------------------------------------------------

# 1-indexed slide numbers that get a "NEEDS UPDATING" banner
NARRATIVE_SLIDES = {
    1, 2, 3, 4, 5, 6,
    9, 10,
    12, 14,
    23, 24, 25,
    37, 38, 39, 40, 41, 42, 43, 44, 45, 46,
    48, 49,
}

# Mapping of slide number -> list of (image_index, chart_func_name, kwargs)
# image_index is the order among image shapes on that slide (0-based, L-to-R)
CHART_SLIDES: dict[int, list[tuple[int, str, dict[str, Any]]]] = {
    11: [(0, "overview_results_chart", {})],
    15: [(0, "composite_satisfaction_chart", {})],
    16: [(0, "question_stacked_bar", {"q_idx": 0, "q_num": 3}),
         (1, "question_yoy_chart", {"q_idx": 0, "q_num": 3})],
    17: [(0, "question_stacked_bar", {"q_idx": 1, "q_num": 4}),
         (1, "question_yoy_chart", {"q_idx": 1, "q_num": 4})],
    18: [(0, "question_stacked_bar", {"q_idx": 2, "q_num": 5}),
         (1, "question_yoy_chart", {"q_idx": 2, "q_num": 5})],
    19: [(0, "question_stacked_bar", {"q_idx": 3, "q_num": 6}),
         (1, "question_yoy_chart", {"q_idx": 3, "q_num": 6})],
    20: [(0, "question_stacked_bar", {"q_idx": 4, "q_num": 7}),
         (1, "question_yoy_chart", {"q_idx": 4, "q_num": 7})],
    21: [(0, "question_stacked_bar", {"q_idx": 5, "q_num": 8}),
         (1, "question_yoy_chart", {"q_idx": 5, "q_num": 8})],
    22: [(0, "question_stacked_bar", {"q_idx": 6, "q_num": 9}),
         (1, "question_yoy_chart", {"q_idx": 6, "q_num": 9})],
    26: [(0, "good_choice_better_serve_chart", {})],
    27: [(0, "tag_frequency_chart", {})],
}

# Tag categories and their slide numbers (slides 28-36)
TAG_CATEGORIES = [
    "Culture/ Virtues",
    "Curriculum",
    "Policies/ Administration",
    "Teachers",
    "Good Outcomes",
    "Community",
    "Facilities",
    "Communication",
    "Extra-curriculars/ Sports",
]

# Build category slide mapping: slides 28-36
for _i, _cat in enumerate(TAG_CATEGORIES):
    slide_num = 28 + _i
    CHART_SLIDES[slide_num] = [
        (0, "category_yoy_chart", {"category": _cat}),
    ]


# ---------------------------------------------------------------------------
# Helper: compute weighted average for a question at a specific level
# ---------------------------------------------------------------------------

def _compute_question_weighted_avg(
    df: pl.DataFrame, question: str, level: str
) -> float:
    """Weighted average (1-4 scale) for one question at one level."""
    scale = QUESTION_SCALES.get(question, [])
    col_name = f"({level}) {question}"
    if col_name not in df.columns or not scale:
        return 0.0

    total_score = 0
    total_count = 0
    for idx, s in enumerate(scale):
        score = len(scale) - idx  # 4, 3, 2, 1
        count = df.filter(pl.col(col_name) == s).height
        total_score += score * count
        total_count += count

    return round(total_score / total_count, 2) if total_count > 0 else 0.0


# ---------------------------------------------------------------------------
# Image replacement
# ---------------------------------------------------------------------------

def _get_image_shapes(slide):
    """Return image shapes on a slide sorted left-to-right by position."""
    from pptx.shapes.picture import Picture
    images = [s for s in slide.shapes if isinstance(s, Picture)]
    images.sort(key=lambda s: (s.left, s.top))
    return images


def _replace_image(slide, image_index: int, new_png_path: str):
    """Replace the Nth image (left-to-right) on a slide with a new PNG."""
    images = _get_image_shapes(slide)
    if image_index >= len(images):
        logger.warning(
            "Slide has %d images but tried to replace index %d",
            len(images), image_index,
        )
        return

    old = images[image_index]
    left, top, width, height = old.left, old.top, old.width, old.height

    # Drop the image relationship before removing the element to avoid
    # orphaned parts that bloat the output and trigger repair warnings.
    from pptx.oxml.ns import qn
    blip = old._element.find('.//' + qn('a:blip'))
    if blip is not None:
        rId = blip.get(qn('r:embed'))
        if rId:
            slide.part.drop_rel(rId)

    sp_tree = slide.shapes._spTree
    sp_tree.remove(old._element)

    # Add new image at same position/size
    slide.shapes.add_picture(new_png_path, left, top, width, height)


# ---------------------------------------------------------------------------
# Text replacement
# ---------------------------------------------------------------------------

def _update_text_by_match(slide, old_text: str, new_text: str) -> bool:
    """Find and replace text in shapes on a slide. Returns True if found."""
    found = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            full_text = "".join(run.text for run in paragraph.runs)
            if old_text in full_text:
                # Replace in the first run, clear the rest
                if paragraph.runs:
                    new_full = full_text.replace(old_text, new_text)
                    paragraph.runs[0].text = new_full
                    for run in paragraph.runs[1:]:
                        run.text = ""
                    found = True
    return found


def _update_table_cell(slide, table_index: int, row: int, col: int, value: str):
    """Update a specific cell in a table on the slide."""
    tables = [s for s in slide.shapes if s.has_table]
    if table_index >= len(tables):
        logger.warning("Slide has %d tables, tried index %d", len(tables), table_index)
        return
    table = tables[table_index].table
    if row < len(table.rows) and col < len(table.columns):
        cell = table.cell(row, col)
        cell.text = str(value)


# ---------------------------------------------------------------------------
# "NEEDS UPDATING" banner
# ---------------------------------------------------------------------------

def _add_needs_updating_banner(slide):
    """Add a large semi-transparent rotated 'NEEDS UPDATING' text box."""
    slide_width = Inches(13.333)  # standard 16:9 width
    slide_height = Inches(7.5)

    box_width = Inches(9)
    box_height = Inches(1.5)
    left = int((slide_width - box_width) / 2)
    top = int((slide_height - box_height) / 2)

    txBox = slide.shapes.add_textbox(left, top, box_width, box_height)

    # Rotate ~25 degrees
    txBox.rotation = -25.0

    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "NEEDS UPDATING"
    font = run.font
    font.size = Pt(72)
    font.bold = True
    font.color.rgb = RGBColor(0xCC, 0x00, 0x00)
    # python-pptx doesn't support text transparency directly via the API,
    # but we can set a lighter color to visually suggest transparency
    font.color.rgb = RGBColor(0xCC, 0x44, 0x44)


# ---------------------------------------------------------------------------
# Chart generation helpers
# ---------------------------------------------------------------------------

def _generate_chart(
    func_name: str,
    kwargs: dict[str, Any],
    year: str,
    df: pl.DataFrame,
    years_data: dict[str, pl.DataFrame],
    tag_data: dict[str, Any],
    free_responses: list[dict[str, Any]],
    tmp_dir: Path,
) -> str | None:
    """Generate a chart PNG and return its path, or None on failure."""
    try:
        if func_name == "question_stacked_bar":
            q_idx = kwargs["q_idx"]
            q_num = kwargs["q_num"]
            question = QUESTIONS[q_idx]
            return generate_question_stacked_bar(df, question, year, tmp_dir, q_num)

        elif func_name == "question_yoy_chart":
            q_idx = kwargs["q_idx"]
            q_num = kwargs["q_num"]
            question = QUESTIONS[q_idx]
            return generate_question_yoy_chart(years_data, question, tmp_dir, q_num)

        elif func_name == "composite_satisfaction_chart":
            paths = generate_composite_satisfaction_chart(years_data, tmp_dir)
            # Returns [png_path, csv_path] — we want the PNG
            return paths[0] if paths else None

        elif func_name == "good_choice_better_serve_chart":
            paths = generate_good_choice_better_serve_chart(
                free_responses, year, tmp_dir
            )
            return paths[0] if paths else None

        elif func_name == "tag_frequency_chart":
            gc_tags = tag_data.get("good_choice", {})
            bs_tags = tag_data.get("better_serve", {})
            if gc_tags or bs_tags:
                return generate_tag_frequency_chart(gc_tags, bs_tags, year, tmp_dir)
            return None

        elif func_name == "category_yoy_chart":
            category = kwargs["category"]
            # years_tag_data format: {year: {tag: count}}
            years_tag_counts = tag_data.get("years_category_counts", {})
            if years_tag_counts:
                return generate_category_yoy_chart(years_tag_counts, category, tmp_dir)
            return None

        elif func_name == "overview_results_chart":
            from app.core.charts import generate_satisfaction_summary_chart
            return generate_satisfaction_summary_chart(df, year, tmp_dir)

        else:
            logger.warning("Unknown chart function: %s", func_name)
            return None

    except Exception:
        logger.exception("Error generating chart %s", func_name)
        return None


# ---------------------------------------------------------------------------
# Data-driven text updates per slide
# ---------------------------------------------------------------------------

def _update_slide_text(
    slide,
    slide_num: int,
    year: str,
    df: pl.DataFrame,
    years_data: dict[str, pl.DataFrame],
    tag_data: dict[str, Any],
    free_responses: list[dict[str, Any]],
):
    """Apply data-driven text replacements for specific slides."""

    # Slides 7-8: participation tables — update response counts
    if slide_num in (7, 8):
        total = df.height
        for level in LEVELS:
            # Count responses for each level
            level_cols = [c for c in df.columns if c.startswith(f"({level})")]
            if level_cols:
                level_count = df.filter(pl.col(level_cols[0]).is_not_null()).height
                _update_text_by_match(slide, f"{level}: ", f"{level}: {level_count}")

    # Slide 13: composite score tables
    elif slide_num == 13:
        overall_avg = compute_overall_average(df)
        _update_text_by_match(slide, "Overall:", f"Overall: {overall_avg:.2f}")
        for level in LEVELS:
            avg = compute_level_average(df, level)
            _update_text_by_match(slide, f"{level}:", f"{level}: {avg:.2f}")

        # Update historical averages for prior years
        for y, y_df in sorted(years_data.items()):
            if y != year:
                y_avg = compute_overall_average(y_df)
                _update_text_by_match(slide, f"{y}:", f"{y}: {y_avg:.2f}")

    # Slides 16-22: per-question historical averages and weighted averages
    elif 16 <= slide_num <= 22:
        q_idx = slide_num - 16  # 0..6
        question = QUESTIONS[q_idx]

        # Update weighted averages for each level
        for level in LEVELS:
            wavg = _compute_question_weighted_avg(df, question, level)
            # Try common annotation patterns
            _update_text_by_match(
                slide, f"{level} Avg", f"{level} Avg: {wavg:.2f}"
            )

        # Update historical year annotations
        for y, y_df in sorted(years_data.items()):
            y_wavg_all = 0.0
            count = 0
            for level in LEVELS:
                v = _compute_question_weighted_avg(y_df, question, level)
                if v > 0:
                    y_wavg_all += v
                    count += 1
            if count:
                avg = round(y_wavg_all / count, 2)
                _update_text_by_match(slide, f"{y}:", f"{y}: {avg:.2f}")

    # Slide 26: Good Choice vs Better Serve counts
    elif slide_num == 26:
        if free_responses:
            respondent_types: dict[str, set[str]] = {}
            for r in free_responses:
                sid = r.get("survey_response_id") or r.get("response_id", "")
                qt = r.get("question_type", "")
                if sid not in respondent_types:
                    respondent_types[sid] = set()
                respondent_types[sid].add(qt)

            total_gc = sum(1 for t in respondent_types.values() if "praise" in t)
            only_gc = sum(
                1 for t in respondent_types.values()
                if "praise" in t and "improvement" not in t
            )
            total_respondents = len(respondent_types)
            only_pct = (
                round(only_gc / total_respondents * 100)
                if total_respondents > 0
                else 0
            )
            # "183" is rendered inside the chart image (already replaced).
            # Match the text pattern "10% ONLY" on the slide.
            _update_text_by_match(slide, "10% ONLY", f"{only_pct}% ONLY")

    # Slides 28-36: tag category response frequency
    elif 28 <= slide_num <= 36:
        cat_idx = slide_num - 28
        if cat_idx < len(TAG_CATEGORIES):
            cat = TAG_CATEGORIES[cat_idx]
            gc_tags = tag_data.get("good_choice", {})
            bs_tags = tag_data.get("better_serve", {})
            freq = gc_tags.get(cat, 0) + bs_tags.get(cat, 0)
            _update_text_by_match(
                slide, "Response Frequency:", f"Response Frequency: {freq}"
            )

    # Slide 47: summary satisfaction rate percentages
    # Template text: "Overall Satisfaction Rate:\n87.5%\nGrammar School: 93.3%\n..."
    # We find the shape containing "Overall Satisfaction Rate:" and replace
    # all percentage patterns with freshly computed values.
    elif slide_num == 47:
        overall_avg = compute_overall_average(df)
        overall_pct = round((overall_avg / 4) * 100, 1)
        level_pcts = {}
        for level in LEVELS:
            avg = compute_level_average(df, level)
            level_pcts[level] = round((avg / 4) * 100, 1)

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full = shape.text_frame.text
            if "Overall Satisfaction Rate" not in full:
                continue
            # Replace percentages in order: overall first, then per-level
            pct_values = [overall_pct] + [level_pcts[lv] for lv in LEVELS]
            pct_iter = iter(pct_values)
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    def _replace_pct(m):
                        try:
                            return f"{next(pct_iter)}%"
                        except StopIteration:
                            return m.group(0)
                    run.text = re.sub(r"\d+\.\d+%", _replace_pct, run.text)
            break


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def generate_deck(
    year: str,
    template_path: str | Path,
    output_dir: str | Path,
    df: pl.DataFrame,
    years_data: dict[str, pl.DataFrame],
    tag_data: dict[str, Any],
    free_responses: list[dict[str, Any]],
) -> Path:
    """Generate a PPTX deck from the template with fresh data.

    Args:
        year: Target survey year (e.g. "2026").
        template_path: Path to the template .pptx file.
        output_dir: Directory for the output .pptx.
        df: DataFrame for the target year.
        years_data: {year: DataFrame} for all years (for YoY charts).
        tag_data: Tag distribution dict with keys "good_choice", "better_serve",
                  and optionally "years_category_counts" for category YoY.
        free_responses: List of free response dicts from PocketBase.

    Returns:
        Path to the generated .pptx file.
    """
    template_path = Path(template_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    prs = Presentation(str(template_path))

    # Generate all charts into a temporary directory
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_path = Path(tmp_dir)

        total_slides = len(prs.slides)
        logger.info("Processing %d slides for year %s", total_slides, year)

        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1  # 1-indexed

            # --- Narrative slides: add banner ---
            if slide_num in NARRATIVE_SLIDES:
                _add_needs_updating_banner(slide)
                logger.debug("Slide %d: added NEEDS UPDATING banner", slide_num)
                continue

            # --- Data-driven chart slides ---
            if slide_num in CHART_SLIDES:
                chart_specs = CHART_SLIDES[slide_num]
                for img_idx, func_name, kwargs in chart_specs:
                    png_path = _generate_chart(
                        func_name, kwargs, year, df, years_data,
                        tag_data, free_responses, tmp_path,
                    )
                    if png_path:
                        _replace_image(slide, img_idx, png_path)
                        logger.debug(
                            "Slide %d: replaced image %d with %s",
                            slide_num, img_idx, func_name,
                        )
                    else:
                        logger.warning(
                            "Slide %d: no chart generated for %s",
                            slide_num, func_name,
                        )

            # --- Data-driven text updates ---
            _update_slide_text(
                slide, slide_num, year, df, years_data,
                tag_data, free_responses,
            )

        # Save inside the tempdir context so temp PNGs are still available
        output_path = output_dir / f"presentation_{year}.pptx"
        prs.save(str(output_path))
        logger.info("Deck saved to %s", output_path)

    return output_path


# ---------------------------------------------------------------------------
# Convenience wrapper that loads data from pipeline/PocketBase
# ---------------------------------------------------------------------------

async def generate_deck_from_pipeline(
    year: str,
    template_path: str | Path,
    years: list[str] | None = None,
) -> Path:
    """High-level helper: loads all data, then calls generate_deck().

    Args:
        year: Target survey year.
        template_path: Path to the template .pptx.
        years: List of years for YoY comparisons. Defaults to all available.

    Returns:
        Path to the generated PPTX file.
    """
    from app.services.pipeline_manager import pipeline_manager
    from app.services.pocketbase_client import pb_client

    settings = get_settings()

    # Determine years
    if not years:
        years = pipeline_manager.get_available_years()
        if year not in years:
            years.append(year)
        years.sort()

    # Load DataFrames for all years
    years_data: dict[str, pl.DataFrame] = {}
    for y in years:
        pipeline = await pipeline_manager.ensure_loaded(y)
        if pipeline.data is None:
            raise ValueError(f"No data loaded for year {y}")
        years_data[y] = pipeline.data

    df = years_data[year]

    # Fetch tag distribution for the target year
    tagging_results = await pb_client.get_full_list(
        "tagging_results", filter_str=f'year = "{year}"'
    )
    free_responses = await pb_client.get_full_list(
        "free_responses", filter_str=f'year = "{year}"'
    )

    good_choice_tags: dict[str, int] = {}
    better_serve_tags: dict[str, int] = {}
    resp_type_map: dict[str, str] = {}

    for fr in free_responses:
        resp_type_map[fr.get("response_id", "")] = fr.get(
            "question_type", "improvement"
        )

    for result in tagging_results:
        rid = result.get("response_id", "")
        qt = resp_type_map.get(rid, "")
        question = result.get("question", "")
        tags = result.get("llm_tags", [])

        if qt == "praise" or "good choice" in question.lower():
            for tag in tags:
                good_choice_tags[tag] = good_choice_tags.get(tag, 0) + 1
        elif qt == "improvement" or "better serve" in question.lower():
            for tag in tags:
                better_serve_tags[tag] = better_serve_tags.get(tag, 0) + 1

    # Build per-year category counts for category YoY charts
    tagging_cache: dict[str, list] = {year: tagging_results}
    years_category_counts: dict[str, dict[str, int]] = {}
    for y in years:
        if y not in tagging_cache:
            tagging_cache[y] = await pb_client.get_full_list(
                "tagging_results", filter_str=f'year = "{y}"'
            )
        cat_counts: dict[str, int] = {}
        for tr in tagging_cache[y]:
            for tag in tr.get("llm_tags", []):
                cat_counts[tag] = cat_counts.get(tag, 0) + 1
        years_category_counts[y] = cat_counts

    tag_data = {
        "good_choice": good_choice_tags,
        "better_serve": better_serve_tags,
        "years_category_counts": years_category_counts,
    }

    output_dir = settings.artifacts_dir / year
    return generate_deck(
        year=year,
        template_path=template_path,
        output_dir=output_dir,
        df=df,
        years_data=years_data,
        tag_data=tag_data,
        free_responses=free_responses,
    )
